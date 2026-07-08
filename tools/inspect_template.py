"""Inspect the IDBI PPT template structure: slides, shapes, placeholders, text."""
from pptx import Presentation
from pptx.util import Emu

SRC = r"C:\Users\ITIK49Q\Downloads\Prototype Submission Deck _ IDBI Innovate.pptx"

prs = Presentation(SRC)
print(f"Slide size: {Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in")
print(f"Total slides: {len(prs.slides)}")
print("=" * 90)

for i, slide in enumerate(prs.slides):
    print(f"\n----- SLIDE {i} (layout: {slide.slide_layout.name}) -----")
    for shape in slide.shapes:
        kind = shape.shape_type
        ph = ""
        if shape.is_placeholder:
            ph = f" [PLACEHOLDER idx={shape.placeholder_format.idx} type={shape.placeholder_format.type}]"
        pos = f"L{Emu(shape.left).inches:.1f} T{Emu(shape.top).inches:.1f} W{Emu(shape.width).inches:.1f} H{Emu(shape.height).inches:.1f}" if shape.left is not None else "no-pos"
        txt = ""
        if shape.has_text_frame:
            txt = " | ".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
        print(f"  [{shape.shape_id}] {kind} name='{shape.name}'{ph} ({pos})")
        if txt:
            print(f"        TEXT: {txt[:160]}")
