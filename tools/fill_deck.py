"""
Fill the IDBI Innovate 2026 PPT template with MSME Financial Health Card content.

Reads the official template, keeps its branding/section headers, and injects
our content (bullets, diagrams, screenshots, live links) into each slide, then
saves a submission-ready .pptx.
"""
from __future__ import annotations

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(__file__)
ROOT = os.path.abspath(os.path.join(HERE, ".."))
ASSETS = os.path.join(ROOT, "deck-assets")
SRC = r"C:\Users\ITIK49Q\Downloads\Prototype Submission Deck _ IDBI Innovate.pptx"
OUT = os.path.join(ROOT, "MSME_Health_Card_IDBI_Innovate_Submission.pptx")

GITHUB = "https://github.com/rtm20/MSME-Health-Card"
LIVE = "https://rtm20-msme-health-card-app-ar3kpv.streamlit.app/"

INK = RGBColor(0x1A, 0x2A, 0x38)
GREEN = RGBColor(0x00, 0x8A, 0x43)
GREY = RGBColor(0x55, 0x66, 0x70)

SLIDE_W = Inches(10)
CONTENT_TOP = Inches(1.55)
CONTENT_LEFT = Inches(0.45)
CONTENT_W = Inches(9.1)
CONTENT_H = Inches(3.65)


def add_bullets(slide, blocks, top=CONTENT_TOP, height=CONTENT_H, size=13):
    """blocks: list of (text, level, bold, color) tuples."""
    tb = slide.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, level, bold, color in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = min(level, 4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size if level else size + 1)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(4)
    return tb


def add_image_fit(slide, path, max_w=9.2, max_h=3.7, top=1.55):
    img = Image.open(path)
    aspect = img.width / img.height
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    left = (10 - w) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))


def two_images(slide, p1, p2, top=1.5, max_h=3.6):
    for i, p in enumerate([p1, p2]):
        img = Image.open(p)
        aspect = img.width / img.height
        w = 4.55
        h = w / aspect
        if h > max_h:
            h = max_h
            w = h * aspect
        left = 0.35 + i * 4.75 + (4.55 - w) / 2
        slide.shapes.add_picture(p, Inches(left), Inches(top + (max_h - h) / 2),
                                 Inches(w), Inches(h))


def set_line(paragraph, label, value, link=None):
    """Rewrite an existing template paragraph as 'label: value' with optional link."""
    for r in list(paragraph.runs):
        r.text = ""
    p = paragraph
    r1 = p.add_run()
    r1.text = f"{label}:  "
    r1.font.bold = True
    r1.font.size = Pt(15)
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = value
    r2.font.size = Pt(15)
    r2.font.color.rgb = GREEN if link else INK
    if link:
        r2.hyperlink.address = link


def main():
    prs = Presentation(SRC)
    s = prs.slides

    # ---- Slide 0: Team details (edit existing box) ----
    box = s[0].shapes[1].text_frame
    # paragraphs: [Team Details, Team name: , Team leader name: , Problem Statement:]
    paras = box.paragraphs
    if len(paras) >= 4:
        set_line(paras[1], "Team name", "[ Add your team name ]")
        set_line(paras[2], "Team leader name", "[ Add team leader name ]")
        set_line(paras[3], "Problem Statement",
                 "PS 3 - Financial Health Score (AI/ML MSME Financial Health Card)")

    # ---- Slide 1: Brief about the idea ----
    add_bullets(s[1], [
        ("An AI/ML-driven MSME Financial Health Card that fuses alternate data "
         "(GST, UPI, Account Aggregator, EPFO) into a single, transparent, "
         "multidimensional Financial Health Score (0-100).", 0, False, INK),
        ("It turns the fragmented digital footprint of New-to-Credit / New-to-Bank "
         "MSMEs into an explainable creditworthiness view - enabling near "
         "real-time, inclusive and prudent lending.", 0, False, INK),
        ("Designed to plug into India's ULI / OCEN / Account Aggregator rails and "
         "into the bank's mobile / loan-origination stack.", 0, False, INK),
        ("Outcome: fewer rejected viable borrowers, faster financial inclusion, "
         "and better portfolio quality.", 0, True, GREEN),
    ], size=14)

    # ---- Slide 2: Opportunities ----
    add_bullets(s[2], [
        ("How it is different", 0, True, GREEN),
        ("Fuses 4 alternate-data sources into ONE score vs siloed point checks", 1, False, INK),
        ("Dual-engine: transparent glass-box pillar score + ML challenger (not a black box)", 1, False, INK),
        ("Per-applicant SHAP explainability - audit-ready for regulated lending", 1, False, INK),
        ("How it solves the problem", 0, True, GREEN),
        ("Scores credit-invisible MSMEs that lack audited financials; expands inclusion", 1, False, INK),
        ("Near real-time re-scoring on fresh consented data", 1, False, INK),
        ("USP", 0, True, GREEN),
        ("Explainable-by-design + interoperable (OCEN/ULI) + inclusion-first, "
         "deployable into the bank's stack", 1, False, INK),
    ], size=12.5)

    # ---- Slide 3: Features ----
    add_bullets(s[3], [
        ("6-pillar Financial Health Score (0-100) with rating band (AAA-C)", 1, False, INK),
        ("Alternate-data aggregation: GST, UPI, Account Aggregator, EPFO", 1, False, INK),
        ("Dual PD: transparent rule-based + XGBoost ML (AUC ~ 0.78, KS ~ 0.49)", 1, False, INK),
        ("Per-applicant SHAP explainability - what drives each borrower's risk", 1, False, INK),
        ("Indicative credit-limit recommendation", 1, False, INK),
        ("What-if simulator for live re-scoring and RM coaching", 1, False, INK),
        ("Portfolio & financial-inclusion analytics + CSV export", 1, False, INK),
        ("OCEN / ULI-ready JSON API with DEPA-AA consent framing", 1, False, INK),
        ("Score-improvement recommendations for the MSME", 1, False, INK),
    ], size=13)

    # ---- Slide 4: Process flow diagram ----
    add_image_fit(s[4], os.path.join(ASSETS, "process_flow.png"), top=1.7)

    # ---- Slide 5: Wireframes / Mock (actual UI) ----
    add_image_fit(s[5], os.path.join(ASSETS, "shot-healthcard.png"), max_h=3.7, top=1.55)

    # ---- Slide 6: Architecture diagram ----
    add_image_fit(s[6], os.path.join(ASSETS, "architecture.png"), top=1.7)

    # ---- Slide 7: Technologies ----
    add_bullets(s[7], [
        ("App & visualisation", 0, True, GREEN),
        ("Python, Streamlit, Plotly", 1, False, INK),
        ("Machine learning & data", 0, True, GREEN),
        ("scikit-learn, XGBoost, native SHAP tree contributions, pandas, NumPy", 1, False, INK),
        ("Scoring engine", 0, True, GREEN),
        ("Transparent weighted 6-pillar model (glass-box, auditable)", 1, False, INK),
        ("Integration & interoperability", 0, True, GREEN),
        ("OCEN / ULI JSON payload, Account Aggregator (DEPA-AA) consent framing", 1, False, INK),
        ("Deployment", 0, True, GREEN),
        ("Streamlit Community Cloud (live), GitHub, Docker-friendly, sandbox-API ready", 1, False, INK),
    ], size=12.5)

    # ---- Slide 8: Estimated implementation cost ----
    add_bullets(s[8], [
        ("MVP / POC: near-zero - open-source stack + free cloud tier (already live)", 0, True, GREEN),
        ("Production (indicative, variable with volume):", 0, True, INK),
        ("Containerised hosting / compute for the scoring microservice", 1, False, INK),
        ("Account Aggregator data-fetch fees (per consent pull)", 1, False, INK),
        ("GST / bureau API costs (per hit)", 1, False, INK),
        ("Model ops, monitoring & audit logging", 1, False, INK),
        ("Marginal cost per assessment is low -> strong unit economics for "
         "inclusion lending at scale.", 0, True, GREEN),
    ], size=13)

    # ---- Slide 9: Snapshots of the prototype ----
    two_images(s[9], os.path.join(ASSETS, "shot-overview.png"),
               os.path.join(ASSETS, "shot-portfolio.png"), top=1.55)

    # ---- Slide 10: Performance / Benchmarking ----
    add_bullets(s[10], [
        ("ML default-risk model: AUC ~ 0.78, KS ~ 0.49  (KS > 0.4 = strong in credit scoring)", 0, True, GREEN),
        ("Demo portfolio: 1,500 MSMEs scored; ~23% New-to-Credit", 1, False, INK),
        ("~70% underwritable (FHS >= 50); 200+ credit-invisible MSMEs unlocked via alt-data", 1, False, INK),
        ("Instant scoring latency; fully reproducible, transparent pipeline", 1, False, INK),
        ("vs baseline structured-data-only default models (~16-22% accuracy per the "
         "bank's problem) -> large uplift PLUS regulator-friendly explainability", 0, True, GREEN),
    ], size=13)

    # ---- Slide 11: Additional / Future development ----
    add_bullets(s[11], [
        ("Live sandbox integrations: Account Aggregator, GST, UPI (NPCI), EPFO", 1, False, INK),
        ("Bureau-data fusion, reject-inference, and bias / fairness monitoring", 1, False, INK),
        ("Sector-specific scorecards and seasonality-aware models", 1, False, INK),
        ("Vernacular RM co-pilot and straight-through auto-decisioning within limits", 1, False, INK),
        ("Direct embed into IDBI mobile app / LOS; ULI / OCEN production onboarding", 1, False, INK),
        ("Continuous re-scoring and early-warning signals for portfolio quality", 1, False, INK),
    ], size=13)

    # ---- Slide 12: Links (edit existing box) ----
    lbox = s[12].shapes[0].text_frame
    lp = lbox.paragraphs
    # paragraphs: [Provide links to your:, GitHub..., Demo Video..., Final Product Link]
    if len(lp) >= 4:
        set_line(lp[1], "GitHub Public Repository", GITHUB, link=GITHUB)
        set_line(lp[2], "Demo Video Link (3 Minutes)", "[ Add your 3-min demo video URL ]")
        set_line(lp[3], "Final Product Link", LIVE, link=LIVE)

    prs.save(OUT)
    print("Saved:", OUT)


if __name__ == "__main__":
    main()
